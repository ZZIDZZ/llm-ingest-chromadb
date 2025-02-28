# utils/extract.py
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file. Skip password-protected files."""
    try:
        reader = PdfReader(file_path)
        if reader.is_encrypted:
            print(f"Skipping password-protected or encrypted PDF: {file_path}")
            return ""
        all_text = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                all_text.append(page_text)
        return "\n".join(all_text)
    # except PdfReadError as e:
        # print(f"Skipping {file_path}. Could not read PDF. Error: {e}")
        # return ""
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_excel(file_path):
    df_dict = pd.read_excel(file_path, sheet_name=None)
    combined_sheets = []
    for _, df in df_dict.items():
        combined_sheets.append("\n".join(df.astype(str).fillna("").values.flatten()))
    return "\n".join(combined_sheets)

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    slides_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)
    return "\n".join(slides_text)

def extract_text(file_path: str) -> str:
    """Determine file type and extract text accordingly."""
    ext = file_path.split(".")[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext == "docx":
        return extract_text_from_docx(file_path)
    elif ext in ["xls", "xlsx"]:
        return extract_text_from_excel(file_path)
    elif ext == "pptx":
        return extract_text_from_pptx(file_path)
    else:
        print(f"Unsupported file type: {file_path}")
        return ""
